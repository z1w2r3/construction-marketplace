"""
PowerPoint 文档解析器模块

解析 .pptx 格式的 PowerPoint 文档
提取幻灯片内容、标题、备注等信息
"""
from typing import Dict, List, Optional, Any
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from .base_parser import BaseParser
from utils import get_logger, ParseError

logger = get_logger(__name__)


class PowerPointParser(BaseParser):
    """PowerPoint 文档解析器"""

    def __init__(self):
        super().__init__()

    def get_supported_extensions(self) -> list:
        """获取支持的文件扩展名"""
        return ['.pptx', '.ppt']

    def parse(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> Dict:
        """
        解析 PowerPoint 文档

        Args:
            file_path: PPT 文档路径
            options: 解析选项
                - max_slides: 最大幻灯片数
                - extract_notes: 是否提取备注

        Returns:
            解析结果字典
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ParseError(
                "缺少 python-pptx 库，请运行: pip install python-pptx"
            )

        options = options or {}
        max_slides = options.get('max_slides', 50)
        extract_notes = options.get('extract_notes', True)

        try:
            # 加载演示文稿
            self.logger.info(f"加载 PowerPoint 文档: {file_path}")
            prs = Presentation(file_path)

            # 提取内容
            content = {}

            # 1. 提取幻灯片
            slides_data = self._extract_slides(prs, max_slides, extract_notes)
            content['slides'] = slides_data

            # 2. 生成摘要
            summary = self._generate_summary(slides_data)

            # 3. 元数据
            metadata = self._extract_metadata(prs)

            return self._create_success_response(
                file_path,
                content,
                summary,
                metadata
            )

        except Exception as e:
            self.logger.error(f"PowerPoint 文档解析失败: {e}", exc_info=True)
            raise ParseError(f"PowerPoint 文档解析失败: {str(e)}")

    def _extract_slides(
        self,
        prs,
        max_slides: int,
        extract_notes: bool
    ) -> List[Dict]:
        """
        提取幻灯片内容

        Args:
            prs: Presentation 对象
            max_slides: 最大幻灯片数
            extract_notes: 是否提取备注

        Returns:
            幻灯片列表
        """
        slides = []

        for i, slide in enumerate(prs.slides):
            if i >= max_slides:
                break

            slide_data = {
                "index": i + 1,
                "title": "",
                "content": [],
                "notes": "",
                "shape_count": len(slide.shapes)
            }

            # 提取形状中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()

                    # 识别标题 (通常是第一个大文本或特定位置的文本)
                    if not slide_data["title"] and (
                        hasattr(shape, 'is_placeholder') and
                        shape.is_placeholder and
                        shape.placeholder_format.type == 1  # 标题占位符
                    ):
                        slide_data["title"] = text
                    else:
                        slide_data["content"].append(text)

            # 提取备注
            if extract_notes and slide.has_notes_slide:
                try:
                    notes_slide = slide.notes_slide
                    notes_text_frame = notes_slide.notes_text_frame
                    if notes_text_frame:
                        slide_data["notes"] = notes_text_frame.text.strip()
                except Exception as e:
                    self.logger.warning(f"提取幻灯片 {i+1} 备注失败: {e}")

            slides.append(slide_data)

        self.logger.info(f"提取幻灯片: {len(slides)} 张")
        return slides

    def _extract_metadata(self, prs) -> Dict:
        """提取元数据"""
        metadata = {}

        try:
            core_props = prs.core_properties

            if hasattr(core_props, 'author') and core_props.author:
                metadata['author'] = core_props.author

            if hasattr(core_props, 'title') and core_props.title:
                metadata['title'] = core_props.title

            if hasattr(core_props, 'created') and core_props.created:
                metadata['created'] = str(core_props.created)

            if hasattr(core_props, 'modified') and core_props.modified:
                metadata['modified'] = str(core_props.modified)

        except Exception as e:
            self.logger.warning(f"提取元数据失败: {e}")

        return metadata

    def _generate_summary(self, slides: List[Dict]) -> Dict:
        """生成摘要"""
        total_content_items = sum(len(slide['content']) for slide in slides)
        slides_with_notes = sum(1 for slide in slides if slide.get('notes'))

        # 提取所有标题
        titles = [slide['title'] for slide in slides if slide['title']]

        summary = {
            "total_slides": len(slides),
            "total_content_items": total_content_items,
            "slides_with_notes": slides_with_notes,
            "slide_titles": titles
        }

        return summary
